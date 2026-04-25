#!/usr/bin/env python3
"""Build sparc-deck.pptx — one-off PPTX builder for SpArc deck."""
from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# Design system
FONT = "Calibri"
TITLE_C = RGBColor(0x1A, 0x1A, 0x1A)
BODY_C = RGBColor(0x37, 0x41, 0x51)
ACCENT_BLUE = RGBColor(0x25, 0x63, 0xEB)
BG = RGBColor(0xF8, 0xF8, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_slide_bg(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def set_notes(slide, text: str) -> None:
    ns = slide.notes_slide
    ns.notes_text_frame.text = text


def add_title(slide, text: str, top: float = 0.45, size: int = 40) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(top), Inches(8.8), Inches(1.1)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = TITLE_C


def add_subtitle(slide, text: str, top: float = 1.25) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(top), Inches(8.8), Inches(0.55)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(20)
    p.font.color.rgb = ACCENT_BLUE


def add_bullets(
    slide, lines: list[str], top: float = 1.85, size: int = 20
) -> None:
    if not lines:
        return
    box = slide.shapes.add_textbox(
        Inches(0.7), Inches(top), Inches(8.6), Inches(4.5)
    )
    tf = box.text_frame
    tf.text = lines[0]
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
    for p in tf.paragraphs:
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = BODY_C
        p.space_after = Pt(6)


def add_flow_strip(slide) -> None:
    """Fund → Click → Attribute → Pay → Verify (blue boxes + arrows)."""
    labels = ["Fund", "Click", "Attribute", "Pay", "Verify"]
    y, h, w = 6.35, 0.42, 1.1
    x0 = 0.55
    gap = 0.12
    for i, lab in enumerate(labels):
        x = x0 + i * (w + gap + 0.15)
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = ACCENT_BLUE
        shp.line.color.rgb = ACCENT_BLUE
        tf = shp.text_frame
        tf.paragraphs[0].text = lab
        tf.paragraphs[0].font.name = FONT
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        if i < len(labels) - 1:
            ax = x + w + 0.02
            t = slide.shapes.add_textbox(
                Inches(ax), Inches(y + 0.1), Inches(0.15), Inches(0.35)
            )
            t.text_frame.text = "→"
            t.text_frame.paragraphs[0].font.size = Pt(14)
            t.text_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE
            t.text_frame.paragraphs[0].font.bold = True


def add_tech_strip(slide) -> None:
    line = "Gemini   |   Circle USDC   |   Arc Testnet   |   Next.js"
    box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.45), Inches(9.0), Inches(0.45)
    )
    p = box.text_frame.paragraphs[0]
    p.text = line
    p.font.name = FONT
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_BLUE
    p.font.bold = True


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    blank = prs.slide_layouts[6]

    # ---- Slide 1: Title
    s1 = prs.slides.add_slide(blank)
    set_slide_bg(s1)
    tbox = s1.shapes.add_textbox(
        Inches(0.5), Inches(1.1), Inches(9.0), Inches(0.9)
    )
    p = tbox.text_frame.paragraphs[0]
    p.text = "SpArc"
    p.font.name = FONT
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = TITLE_C
    add_subtitle(
        s1,
        "A campaign marketplace engine powered by real engagement.",
        top=1.95,
    )
    add_bullets(
        s1,
        [
            "Fund campaigns in USDC.",
            "Generate creator-facing briefs with Gemini.",
            "Watch clicks resolve into payouts on-chain.",
        ],
        top=2.65,
        size=19,
    )
    set_notes(
        s1,
        "This is SpArc — a campaign funding platform that helps managers "
        "create campaign briefs and pays creators automatically based on "
        "real user behavior.",
    )

    # ---- Slide 2: Problem
    s2 = prs.slides.add_slide(blank)
    set_slide_bg(s2)
    add_title(s2, "Campaign marketing is still rigid and hard to trust", top=0.4, size=30)
    add_bullets(
        s2,
        [
            "Brands pay upfront without clear performance feedback.",
            "Creator campaigns are expensive and inflexible.",
            "Attribution is often incomplete or hidden.",
            "ROI is hard to verify.",
        ],
        top=1.35,
    )
    set_notes(
        s2,
        "Traditional creator campaigns are expensive to launch, hard to measure, and slow to pay out.",
    )

    # ---- Slide 3: Solution
    s3 = prs.slides.add_slide(blank)
    set_slide_bg(s3)
    add_title(s3, "A campaign wallet that pays for real behavior", top=0.4, size=32)
    add_bullets(
        s3,
        [
            "Campaigns are funded in USDC.",
            "Businesses can create and manage campaigns.",
            "Approved creators participate in the campaign.",
            "Clicks create instant micropayments.",
            "Conversions trigger attribution-based payouts.",
            "Every payout is verifiable on-chain.",
        ],
        top=1.25,
    )
    add_flow_strip(s3)
    set_notes(
        s3,
        "We replace rigid influencer contracts with a performance-based campaign wallet.",
    )

    # ---- Slide 4: Product
    s4 = prs.slides.add_slide(blank)
    set_slide_bg(s4)
    add_title(s4, "SpArc gives managers a live campaign workspace", top=0.4, size=30)
    add_bullets(
        s4,
        [
            "Generate campaign briefs with Gemini.",
            "Track click events and micropayments.",
            "View attribution splits across creators.",
            "See payout execution and txHash proof.",
        ],
        top=1.35,
    )
    set_notes(
        s4,
        "This is the working system. The dashboard shows the causal chain from user behavior to payout in real time.",
    )

    # ---- Slide 5: Closing
    s5 = prs.slides.add_slide(blank)
    set_slide_bg(s5)
    add_title(s5, "Campaign marketing, redefined", top=0.5, size=40)
    add_bullets(
        s5,
        [
            "More flexible for brands.",
            "More accessible for creators.",
            "More transparent for everyone.",
            "Built for real-time, on-chain proof.",
        ],
        top=1.45,
    )
    set_notes(
        s5,
        "SpArc makes campaign marketing more democratic: anyone can fund a campaign, creators can join, and payment follows real performance.",
    )

    # ---- Slide 6: Tech note
    s6 = prs.slides.add_slide(blank)
    set_slide_bg(s6)
    add_title(s6, "How it works:", top=0.45, size=40)
    add_bullets(
        s6,
        [
            "Gemini generates campaign-facing copy.",
            "Circle handles the USDC payment flow.",
            "Attribution and payouts remain deterministic.",
            "On-chain proof is live and verifiable.",
        ],
        top=1.3,
    )
    add_tech_strip(s6)
    set_notes(
        s6,
        "This is intentionally scoped for a reliable hackathon demo.",
    )

    # ---- Slide 7: Why this model
    s7 = prs.slides.add_slide(blank)
    set_slide_bg(s7)
    add_title(s7, "Why this model matters", top=0.5, size=40)
    add_bullets(
        s7,
        [
            "Lowers the barrier to entry for businesses",
            "Lowers the barrier to entry for creators",
            "Creates a more open marketplace",
            "Shifts the model from fixed sponsorships to accessible, behavior-based participation",
        ],
        top=1.4,
    )

    # ---- Slide 8: Live demo
    s8 = prs.slides.add_slide(blank)
    set_slide_bg(s8)
    add_title(s8, "Live demo: click → attribution → instant payout", top=0.4, size=32)
    add_bullets(
        s8,
        [
            "Open a funded campaign.",
            "Trigger creator clicks.",
            "Watch the attribution update.",
            "Confirm payout execution on-chain.",
        ],
        top=1.35,
    )
    set_notes(
        s8,
        "We'll show one campaign from start to finish so you can see the full payment flow in action.",
    )

    import os
    from pathlib import Path

    preferred = "/mnt/user-data/outputs/sparc-deck.pptx"
    fallback = (
        Path(__file__).resolve().parent.parent / "user-data" / "outputs" / "sparc-deck.pptx"
    )
    out = preferred
    try:
        os.makedirs(os.path.dirname(out), exist_ok=True)
        prs.save(out)
    except OSError:
        out = str(fallback)
        fallback.parent.mkdir(parents=True, exist_ok=True)
        prs.save(out)
        print(
            f"Note: {preferred} is not writable here; used workspace path instead."
        )
    assert len(prs.slides) == 8
    print(f"Wrote {out} with {len(prs.slides)} slides.")


if __name__ == "__main__":
    main()
